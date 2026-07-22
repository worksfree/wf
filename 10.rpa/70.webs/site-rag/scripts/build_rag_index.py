# -*- coding: utf-8 -*-
"""rag_sources.txt 화이트리스트의 문서를 읽어 Cloudflare Vectorize 업로드용 NDJSON(vectors.ndjson)을 생성한다.

- 임베딩: Ollama 로컬 서버 (모델: bge-m3, 1024차원) — 비용 0원, 완전 오프라인
  질의 시점(Cloudflare Worker → PC 터널)에도 동일 모델을 동일 방식(Ollama API)으로 호출해
  인덱스 빌드와 질의 임베딩의 모델 불일치를 방지한다.
- 소스: cards/*.json (biz-support 지원제도 카드), *.md (NAS 가이드/블로그), *.txt/*.pdf/*.pptx (일반 문서)
  PDF/PPTX 파싱은 pypdf/python-pptx 필요 (requirements.txt). 스캔본 PDF 등 텍스트 추출이 안 되는 파일은
  건너뛴다.
- 자료 추가 시 rag_sources.txt 에 경로를 한 줄 추가한 뒤 본 스크립트 재실행

출력은 정적 JSON(rag_index.json)이 아니라 `wrangler vectorize insert`용 NDJSON이다 — 청크 수가
수천 건으로 늘면(45.Slife 포함 시 6,305개, 80MB) Worker가 인덱스 전체를 fetch+parse하는 방식은
Cloudflare Worker 메모리 한도(128MB)에 위험할 정도로 근접해, 검색 자체를 Cloudflare Vectorize에
맡기는 구조로 전환했다. 벡터 id는 원본 청크 id(한글·특수문자 포함)를 그대로 쓰지 않고 sha1 해시로
치환한다 — Vectorize id 허용 문자·길이 제약을 안전하게 우회하기 위함이며, 원본 식별자는
metadata.chunk_ref 에 보존한다.
"""
import hashlib
import io
import json
import re
import sys
import time
import urllib.error
import urllib.request
from pathlib import Path

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8")
sys.stdout.reconfigure(line_buffering=True)  # 줄 단위 즉시 flush — 진행 상황이 실시간으로 보이도록

BASE = Path(__file__).resolve().parent.parent
SOURCES_FILE = BASE / "rag_sources.txt"
VECTORS_PATH = BASE / "vectors.ndjson"
OLLAMA_URL = "http://localhost:11434/api/embeddings"
MAX_METADATA_TEXT_CHARS = 1800  # Vectorize 메타데이터 크기 한도 방어용 상한
EMBED_MODEL = "bge-m3"
CHUNK_MAX_CHARS = 1000
CHUNK_OVERLAP = 150


def load_source_patterns() -> list[str]:
    patterns = []
    for line in SOURCES_FILE.read_text(encoding="utf-8").splitlines():
        line = line.split("#", 1)[0].strip()
        if line:
            patterns.append(line)
    return patterns


def resolve_files(patterns: list[str]) -> list[Path]:
    files = []
    for pat in patterns:
        files.extend(sorted(BASE.glob(pat)))
    return files


def chunk_text(text: str, max_chars: int = CHUNK_MAX_CHARS, overlap: int = CHUNK_OVERLAP) -> list[str]:
    """문단 경계를 우선 존중하는 슬라이딩 윈도우 청킹."""
    text = text.strip()
    if len(text) <= max_chars:
        return [text] if text else []
    paras = [p.strip() for p in re.split(r"\n{2,}", text) if p.strip()]
    chunks, buf = [], ""
    for p in paras:
        if len(buf) + len(p) + 2 <= max_chars:
            buf = f"{buf}\n\n{p}" if buf else p
        else:
            if buf:
                chunks.append(buf)
            buf = (buf[-overlap:] + "\n\n" + p) if buf and len(p) < max_chars else p
    if buf:
        chunks.append(buf)
    # 문단이 하나도 max_chars 이하로 안 쪼개지면(예: 표) 강제 슬라이딩 윈도우로 재분할
    final = []
    for c in chunks:
        if len(c) <= max_chars * 1.3:
            final.append(c)
        else:
            i = 0
            while i < len(c):
                final.append(c[i:i + max_chars])
                i += max_chars - overlap
    return final


def chunks_from_markdown(path: Path) -> list[dict]:
    raw = path.read_text(encoding="utf-8")
    h1_match = re.search(r"^#\s+(.+)$", raw, re.MULTILINE)
    doc_title = h1_match.group(1).strip() if h1_match else path.stem

    # ## 헤딩 기준 섹션 분할 (없으면 문서 전체를 하나의 섹션으로 취급)
    sections = re.split(r"^##\s+(.+)$", raw, flags=re.MULTILINE)
    out = []
    if len(sections) == 1:
        for i, ch in enumerate(chunk_text(raw)):
            out.append({"title": doc_title, "text": ch, "chunk_no": i})
    else:
        # re.split with a capturing group returns [pre, heading1, body1, heading2, body2, ...]
        pre = sections[0].strip()
        if pre:
            for i, ch in enumerate(chunk_text(pre)):
                out.append({"title": doc_title, "text": ch, "chunk_no": i})
        for j in range(1, len(sections), 2):
            heading = sections[j].strip()
            body = sections[j + 1] if j + 1 < len(sections) else ""
            section_title = f"{doc_title} — {heading}"
            for i, ch in enumerate(chunk_text(body)):
                out.append({"title": section_title, "text": ch, "chunk_no": i})
    return out


def chunks_from_txt(path: Path) -> list[dict]:
    raw = path.read_text(encoding="utf-8", errors="ignore")
    return [{"title": path.stem, "text": ch, "chunk_no": i} for i, ch in enumerate(chunk_text(raw))]


def chunks_from_pdf(path: Path) -> list[dict]:
    from pypdf import PdfReader

    reader = PdfReader(str(path))
    pages_text = []
    for page in reader.pages:
        try:
            pages_text.append(page.extract_text() or "")
        except Exception:
            pages_text.append("")
    raw = "\n\n".join(t for t in pages_text if t.strip())
    if not raw.strip():
        return []  # 스캔본 등 텍스트 추출 불가 PDF — 건너뜀
    return [{"title": path.stem, "text": ch, "chunk_no": i} for i, ch in enumerate(chunk_text(raw))]


def chunks_from_pptx(path: Path) -> list[dict]:
    from pptx import Presentation

    prs = Presentation(str(path))
    slide_texts = []
    for slide in prs.slides:
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                t = shape.text_frame.text.strip()
                if t:
                    texts.append(t)
        if texts:
            slide_texts.append("\n".join(texts))
    raw = "\n\n".join(slide_texts)
    if not raw.strip():
        return []
    return [{"title": path.stem, "text": ch, "chunk_no": i} for i, ch in enumerate(chunk_text(raw))]


def chunks_from_card_json(path: Path) -> list[dict]:
    cards = json.loads(path.read_text(encoding="utf-8"))
    out = []
    for c in cards:
        text = "\n".join(filter(None, [
            c.get("name"),
            "카테고리: " + c.get("category", "") if c.get("category") else None,
            "대상: " + c.get("target", "") if c.get("target") else None,
            "혜택: " + c.get("benefit", "") if c.get("benefit") else None,
            "방법: " + c.get("how", "") if c.get("how") else None,
            "시기: " + c.get("when", "") if c.get("when") else None,
            c.get("summary"),
        ]))
        out.append({
            "title": c.get("name", c.get("id", "")),
            "text": text,
            "chunk_no": 0,
            "card_id": c.get("id"),
        })
    return out


def ollama_embed(text: str) -> list[float]:
    body = json.dumps({"model": EMBED_MODEL, "prompt": text}).encode("utf-8")
    req = urllib.request.Request(
        OLLAMA_URL, data=body, headers={"Content-Type": "application/json"}, method="POST"
    )
    with urllib.request.urlopen(req, timeout=60) as resp:
        data = json.loads(resp.read().decode("utf-8"))
    return data["embedding"]


def main():
    patterns = load_source_patterns()
    files = resolve_files(patterns)
    print(f"소스 파일 {len(files)}개 발견 (rag_sources.txt 화이트리스트 기준)")

    LOADERS = {
        ".json": chunks_from_card_json,
        ".md": chunks_from_markdown,
        ".txt": chunks_from_txt,
        ".pdf": chunks_from_pdf,
        ".pptx": chunks_from_pptx,
    }

    raw_chunks = []
    for fi, f in enumerate(files):
        loader = LOADERS.get(f.suffix.lower())
        if not loader:
            continue
        t0 = time.monotonic()
        print(f"  [{fi+1}/{len(files)}] 로딩: {f.name} ({f.stat().st_size/1024:.0f} KB)")
        try:
            # relative_to(BASE.parent) 실패 시(예: 드라이브 루트 밖 경로) 절대경로로 대체
            try:
                rel = f.relative_to(BASE.parent).as_posix()
            except ValueError:
                rel = str(f)
            items = loader(f)
        except Exception as e:
            print(f"    로드 실패: {e} — 건너뜀")
            continue
        elapsed = time.monotonic() - t0
        if elapsed > 5:
            print(f"    ⚠ {elapsed:.1f}초 소요 (느림)")
        for it in items:
            raw_chunks.append({**it, "source": rel})
    print(f"청크 {len(raw_chunks)}개 생성 완료 — 임베딩 시작 (Ollama: {EMBED_MODEL})")

    try:
        urllib.request.urlopen("http://localhost:11434/", timeout=3)
    except Exception:
        print("⚠ Ollama(localhost:11434)에 연결할 수 없습니다. `ollama serve` 실행 상태를 확인하세요.")
        sys.exit(1)

    written = 0
    dim = 0
    with VECTORS_PATH.open("w", encoding="utf-8") as out:
        for i, rc in enumerate(raw_chunks):
            cid = f"{rc.get('card_id') or Path(rc['source']).stem}::{rc['chunk_no']}"
            try:
                emb = ollama_embed(rc["text"])
            except (urllib.error.URLError, KeyError) as e:
                print(f"  [{i+1}/{len(raw_chunks)}] 임베딩 실패 ({cid}): {e} — 건너뜀")
                continue
            vec_id = hashlib.sha1(cid.encode("utf-8")).hexdigest()
            text = rc["text"]
            if len(text) > MAX_METADATA_TEXT_CHARS:
                text = text[:MAX_METADATA_TEXT_CHARS] + "…"
            record = {
                "id": vec_id,
                "values": [round(float(x), 6) for x in emb],
                "metadata": {
                    "title": rc["title"],
                    "source": rc["source"],
                    "text": text,
                    "chunk_ref": cid,
                },
            }
            out.write(json.dumps(record, ensure_ascii=False) + "\n")
            dim = len(record["values"])
            written += 1
            if (i + 1) % 20 == 0 or i + 1 == len(raw_chunks):
                print(f"  [{i+1}/{len(raw_chunks)}] 임베딩 완료")

    print(f"NDJSON 저장: {VECTORS_PATH} ({VECTORS_PATH.stat().st_size/1024:.0f} KB, "
          f"벡터 {written}개, dim={dim})")
    print("업로드: npx wrangler vectorize insert biz-rag-index --file=vectors.ndjson "
          "(synology-web/workers/biz-rag 폴더에서 상대경로 조정)")


if __name__ == "__main__":
    main()
